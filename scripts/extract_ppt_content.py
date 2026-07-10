"""Extract all text content from the 4 CoSell PPTs for review."""
import json
import os
from pptx import Presentation
from pptx.util import Emu

SRC_DIR = r"C:\Users\v-adevashish\OneDrive - Microsoft\Desktop\SupportAgentDomain"
OUT_JSON = r"c:\WorkFAST-main\generated-content\support-agent-domain\ppt_extract.json"
OUT_TXT = r"c:\WorkFAST-main\generated-content\support-agent-domain\ppt_extract.txt"

FILES = [
    "CoSell 301.pptx",
    "CoSell_Domain_Business_Logic_Guide.pptx",
    "CRM Cosell Deep dive.pptx",
    "GPSCRM Co-sell.pptx",
]


def iter_shapes(shapes):
    """Recursively yield shapes, descending into group shapes."""
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def extract_table(shape):
    rows = []
    tbl = shape.table
    for r in tbl.rows:
        cells = []
        for c in r.cells:
            cells.append(c.text.strip())
        rows.append(cells)
    return rows


def extract_slide(slide, idx):
    data = {"slide": idx, "title": None, "texts": [], "tables": [], "notes": None}
    for shape in iter_shapes(slide.shapes):
        # Title detection
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                is_title = False
                try:
                    if shape == slide.shapes.title:
                        is_title = True
                except Exception:
                    pass
                if is_title and data["title"] is None:
                    data["title"] = txt
                else:
                    data["texts"].append(txt)
        if shape.has_table:
            data["tables"].append(extract_table(shape))
    # Notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            data["notes"] = notes
    return data


def main():
    os.makedirs(os.path.dirname(OUT_JSON), exist_ok=True)
    all_data = {}
    txt_lines = []
    for fname in FILES:
        path = os.path.join(SRC_DIR, fname)
        if not os.path.exists(path):
            txt_lines.append(f"!!! MISSING FILE: {fname}")
            continue
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            slides.append(extract_slide(slide, i))
        all_data[fname] = slides

        txt_lines.append("=" * 100)
        txt_lines.append(f"FILE: {fname}  ({len(slides)} slides)")
        txt_lines.append("=" * 100)
        for s in slides:
            txt_lines.append(f"\n--- Slide {s['slide']} ---")
            if s["title"]:
                txt_lines.append(f"TITLE: {s['title']}")
            for t in s["texts"]:
                txt_lines.append(f"TEXT:\n{t}")
            for ti, tbl in enumerate(s["tables"], start=1):
                txt_lines.append(f"TABLE {ti}:")
                for row in tbl:
                    txt_lines.append("  | " + " | ".join(row))
            if s["notes"]:
                txt_lines.append(f"NOTES:\n{s['notes']}")

    with open(OUT_JSON, "w", encoding="utf-8") as f:
        json.dump(all_data, f, indent=2, ensure_ascii=False)
    with open(OUT_TXT, "w", encoding="utf-8") as f:
        f.write("\n".join(txt_lines))

    # Summary
    print("Extraction complete.")
    for fname, slides in all_data.items():
        n_tables = sum(len(s["tables"]) for s in slides)
        n_notes = sum(1 for s in slides if s["notes"])
        print(f"  {fname}: {len(slides)} slides, {n_tables} tables, {n_notes} notes")
    print(f"\nJSON: {OUT_JSON}")
    print(f"TXT:  {OUT_TXT}")


if __name__ == "__main__":
    main()
